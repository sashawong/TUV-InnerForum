# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
SCREENSHOT_DIR = BASE / "screenshots"
PPT_PATH = BASE / "TUVforum_2026-04-02_update_summary_fixed.pptx"

GREEN = RGBColor(22, 119, 52)
LIGHT_GREEN = RGBColor(234, 246, 237)
DARK = RGBColor(34, 34, 34)
MUTED = RGBColor(95, 99, 104)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.9))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = DARK

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5))
        p2 = tx2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.9, top=1.7, width=11.5, height=4.8, font_size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.bullet = True


def add_tag(slide, text, left, top, width=1.5, height=0.42, fill=LIGHT_GREEN, color=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color


def add_image(slide, path, left, top, width=None, height=None):
    if path.exists():
        slide.shapes.add_picture(
            str(path),
            Inches(left),
            Inches(top),
            width=Inches(width) if width else None,
            height=Inches(height) if height else None,
        )


slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()
add_title(slide, "TUV Inner Forum 今日更新汇总", "日期：2026-04-02 | 汇总今天已完成的功能、上线状态与 AI 能力调整")
add_tag(slide, "已部署上线", 0.65, 1.55, width=1.55)
add_tag(slide, "含页面截图", 2.35, 1.55, width=1.55)
add_tag(slide, "AI功能更新", 4.05, 1.55, width=1.7)
add_bullets(
    slide,
    [
        "论坛已完成搜索优化、回复编辑/删除、积分体系、积分榜游戏化展示、AI 助手、消息提醒等一整套更新。",
        "今天的重点收口在 AI 法规助手：修正接口理解、增加“正在生成中”状态，并把报告能力改成“汇总最近 N 篇帖子后生成新帖”。",
        "代码已同步到 GitHub，并部署到云服务器运行中的前后端容器。",
    ],
    top=2.1,
    height=3.9,
    font_size=20,
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "一、今天完成的主要功能")
add_bullets(
    slide,
    [
        "搜索体验升级：从简单过滤改成弹窗结果页，并按标题命中、正文命中、回复命中分列展示。",
        "回复管理补齐：用户可以删除自己的回复，也可以重新编辑自己的回复。",
        "积分体系上线：新增用户积分属性，按发帖与回复自动累计，并增加积分排行榜。",
        "排行榜视觉优化：按“青铜、白银、黄金、铂金、钻石、星耀、王者”等段位分层展示，并附带小图标。",
        "管理员不参与积分榜统计，避免干扰真实用户排名。",
    ],
    top=1.65,
    height=5.2,
    font_size=19,
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "二、AI 法规助手与提醒机制")
add_bullets(
    slide,
    [
        "创建“AI法规助手”用户，用于生成论坛回复与汇总内容，不参与积分榜。",
        "帖子中可点击“邀请AI回答”，AI 会结合当前帖子与现有回复输出一条新的建议回复。",
        "新增“正在生成中”状态，用户在等待 AI 回复时可以看到明确提示。",
        "左侧新增“AI助手”页面，用于统一生成汇总帖与配置模型接入。",
        "新增新帖提醒与回复消息提醒，用户可在顶部通知入口查看并标记已读。",
    ],
    top=1.65,
    height=5.2,
    font_size=19,
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "三、AI 汇总逻辑调整")
add_bullets(
    slide,
    [
        "原先的设计是：针对“当前单个帖子”生成一篇报告贴。",
        "现在已改成：汇总最近 N 篇帖子及其回复，生成一篇新的论坛总结帖。N 为可选数字，默认值为 10。",
        "这样更符合“周报/阶段总结”的使用场景，也更适合把零散讨论沉淀成法规信息汇总。",
        "帖子详情页保留“邀请AI回答”；汇总生成入口统一收敛到 AI 助手页面，避免用户理解混乱。",
        "接口层已支持按数量汇总，前端页面也支持直接输入最近汇总篇数。",
    ],
    top=1.65,
    height=5.3,
    font_size=18,
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "四、页面截图")
add_image(slide, SCREENSHOT_DIR / "home-final.png", 0.45, 1.5, width=4.1)
add_image(slide, SCREENSHOT_DIR / "leaderboard-final.png", 4.62, 1.5, width=4.1)
add_image(slide, SCREENSHOT_DIR / "topic-detail-final.png", 8.79, 1.5, width=4.1)
for text, left in [("首页", 1.9), ("积分榜", 6.1), ("帖子详情页", 9.95)]:
    cap = slide.shapes.add_textbox(Inches(left), Inches(6.25), Inches(1.8), Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = DARK
note = slide.shapes.add_textbox(Inches(0.8), Inches(6.65), Inches(11.8), Inches(0.4))
p = note.text_frame.paragraphs[0]
r = p.add_run()
r.text = "截图来自当前线上环境，已登录后台后抓取，用于展示首页、积分榜与帖子详情页的实际效果。"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(11)
r.font.color.rgb = MUTED

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "五、部署与验证结果")
add_bullets(
    slide,
    [
        "代码已同步到 GitHub，并已部署到云服务器目录 /www/wwwroot/TUV-InnerForum。",
        "前端与后端均通过 Docker 重建，论坛页面在线可访问。",
        "前端健康检查返回 200 OK；后端排行榜接口也返回正常数据。",
        "数据库中已补充中文帖子与回复示例，积分榜与搜索功能可以直接演示。",
        "OpenRouter 接口 404 的问题已说明：该地址需要程序以 POST 调用，浏览器直接 GET 打开会返回 Not Found。",
    ],
    top=1.7,
    height=5.2,
    font_size=19,
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "六、后续建议")
add_bullets(
    slide,
    [
        "给 AI 汇总增加筛选条件，例如“只汇总有回复的帖子”或“按标签汇总”。",
        "把 AI 回复和汇总过程做成更明显的按钮 loading、进度提示或完成后自动定位。",
        "继续优化前端体积，当前生产包已超过 500 kB，可考虑分包与懒加载。",
        "如果后续切换豆包模型，可直接复用现有 OpenAI 兼容接入方式，只需要替换接口地址、模型名和 API Key。",
    ],
    top=1.7,
    height=5.2,
    font_size=19,
)

prs.save(PPT_PATH)
print(PPT_PATH)
